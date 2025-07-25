from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Blank slide
slide = prs.slides.add_slide(slide_layout)

# Set slide dimensions (16:9)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(12.33), Inches(0.4))
title_frame = title_box.text_frame
title_frame.text = "Royal Gazette e-Publishing System Architecture"
title_frame.paragraphs[0].font.size = Pt(26)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Define colors - matching the screenshot
colors = {
    'user': RGBColor(255, 152, 0),  # Orange
    'frontend': RGBColor(76, 175, 80),  # Green
    'backend': RGBColor(244, 67, 54),  # Red
    'database': RGBColor(33, 150, 243),  # Blue
    'security': RGBColor(156, 39, 176),  # Purple
    'api': RGBColor(96, 125, 139),  # Blue-gray
    'backup': RGBColor(158, 158, 158),  # Gray
    'layer_bg': RGBColor(250, 250, 250),  # Light gray background
    'layer_border': RGBColor(224, 224, 224),  # Light gray border
}

# Layer dimensions - INCREASED
layer_x = 0.4
layer_width = 9.3
layer_height = 1.3  # Increased from 1.2
layer_spacing = 0.15  # Increased spacing

# Component dimensions - INCREASED
comp_height = 0.5  # Increased from 0.6
comp_width = 1.9  # Increased from 1.5
comp_spacing = 0.1

# 4. User Layer (TOP)
user_y = 0.55
user_layer = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, 
    Inches(layer_x), Inches(user_y), 
    Inches(layer_width), Inches(layer_height)
)
user_layer.fill.solid()
user_layer.fill.fore_color.rgb = colors['layer_bg']
user_layer.line.color.rgb = colors['layer_border']
user_layer.line.width = Pt(1.5)

# User Layer Label
user_label = slide.shapes.add_textbox(Inches(layer_x + 0.2), Inches(user_y + 0.05), Inches(2), Inches(0.3))
user_label.text_frame.text = "4. User Layer"
user_label.text_frame.paragraphs[0].font.size = Pt(14)
user_label.text_frame.paragraphs[0].font.bold = True

# User Layer Components - First row
user_components_1 = [
    {"text": "Government Agencies", "x": 1.0},
    {"text": "Local Government", "x": 3.1},
    {"text": "Gazette Staff", "x": 5.2},
    {"text": "Printing House", "x": 7.3},
]

for comp in user_components_1:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(user_y + 0.35),
        Inches(comp_width), Inches(comp_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['user']
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# User Layer Components - Second row
user_components_2 = [
    {"text": "System Admin", "x": 3.5, "color": RGBColor(211, 47, 47)},
    {"text": "Public Users", "x": 5.6},
]

for comp in user_components_2:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(user_y + 0.95),
        Inches(comp_width), Inches(comp_height - 0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = comp.get("color", colors['user'])
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Update todo status
# Marking task 1 as in_progress - Increase layer container heights

# 3. Frontend Layer
frontend_y = user_y + layer_height + 0.4
frontend_layer = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(layer_x), Inches(frontend_y),
    Inches(layer_width), Inches(layer_height)
)
frontend_layer.fill.solid()
frontend_layer.fill.fore_color.rgb = colors['layer_bg']
frontend_layer.line.color.rgb = colors['layer_border']
frontend_layer.line.width = Pt(1.5)

# Frontend Layer Label
frontend_label = slide.shapes.add_textbox(Inches(layer_x + 0.2), Inches(frontend_y + 0.05), Inches(2), Inches(0.3))
frontend_label.text_frame.text = "3. Frontend Layer"
frontend_label.text_frame.paragraphs[0].font.size = Pt(14)
frontend_label.text_frame.paragraphs[0].font.bold = True

# Frontend Components - First row with larger boxes
frontend_components_1 = [
    {"text": "Internal Portal", "x": 0.8, "w": 2.5},
    {"text": "Public Portal", "x": 3.5, "w": 2.5},
    {"text": "Responsive Design", "x": 6.2, "w": 2.3, "color": RGBColor(0, 150, 136)},
]

for comp in frontend_components_1:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(frontend_y + 0.35),
        Inches(comp.get("w", 2.2)), Inches(comp_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = comp.get("color", colors['frontend'])
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Frontend Components - Second row
frontend_components_2 = [
    {"text": "Submission Forms", "x": 1.0},
    {"text": "Search Interface", "x": 3.1},
    {"text": "Document Viewer", "x": 5.2},
    {"text": "Multi-tab Support", "x": 7.3},
]

for comp in frontend_components_2:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(frontend_y + 0.95),
        Inches(comp_width), Inches(comp_height - 0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['frontend']
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# 2. Backend Layer
backend_y = frontend_y + layer_height + 0.4
backend_layer = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(layer_x), Inches(backend_y),
    Inches(layer_width), Inches(layer_height)
)
backend_layer.fill.solid()
backend_layer.fill.fore_color.rgb = colors['layer_bg']
backend_layer.line.color.rgb = colors['layer_border']
backend_layer.line.width = Pt(1.5)

# Backend Layer Label
backend_label = slide.shapes.add_textbox(Inches(layer_x + 0.2), Inches(backend_y + 0.05), Inches(2), Inches(0.3))
backend_label.text_frame.text = "2. Backend Layer"
backend_label.text_frame.paragraphs[0].font.size = Pt(14)
backend_label.text_frame.paragraphs[0].font.bold = True

# Backend Components - First row
backend_components_1 = [
    {"text": "Core Services", "x": 1.0},
    {"text": "Document Processing", "x": 3.1},
    {"text": "Workflow Engine", "x": 5.2},
    {"text": "Template Engine", "x": 7.3},
]

for comp in backend_components_1:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(backend_y + 0.35),
        Inches(comp_width), Inches(comp_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['backend']
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Backend Components - Second row
backend_components_2 = [
    {"text": "Security Module", "x": 1.0, "color": colors['security']},
    {"text": "Search Service", "x": 3.1},
    {"text": "Authentication", "x": 5.2, "color": colors['security']},
    {"text": "Activity Logger", "x": 7.3},
]

for comp in backend_components_2:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(backend_y + 0.95),
        Inches(comp_width), Inches(comp_height - 0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = comp.get("color", colors['backend'])
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# 1. Database Layer (BOTTOM)
database_y = backend_y + layer_height + 0.4
database_layer = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(layer_x), Inches(database_y),
    Inches(layer_width), Inches(layer_height - 0.3)
)
database_layer.fill.solid()
database_layer.fill.fore_color.rgb = colors['layer_bg']
database_layer.line.color.rgb = colors['layer_border']
database_layer.line.width = Pt(1.5)

# Database Layer Label
db_label = slide.shapes.add_textbox(Inches(layer_x + 0.2), Inches(database_y + 0.05), Inches(2), Inches(0.3))
db_label.text_frame.text = "1. Database Layer"
db_label.text_frame.paragraphs[0].font.size = Pt(14)
db_label.text_frame.paragraphs[0].font.bold = True

# Database Components - Single row with better spacing
db_components = [
    {"text": "Main Database", "x": 0.8},
    {"text": "Document Storage", "x": 2.8},
    {"text": "Activity Logs DB", "x": 4.8},
    {"text": "Search Index", "x": 6.8},
]

for comp in db_components:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(comp["x"]), Inches(database_y + 0.35),
        Inches(1.8), Inches(comp_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['database']
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Backup Systems - positioned to the right
backup_shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.7), Inches(database_y + 0.35),
    Inches(0.9), Inches(comp_height)
)
backup_shape.fill.solid()
backup_shape.fill.fore_color.rgb = colors['backup']
backup_shape.text_frame.text = "Backup"
backup_shape.text_frame.paragraphs[0].font.size = Pt(10)
backup_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
backup_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
backup_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# 5. API Layer (RIGHT SIDE) - Adjusted height
api_layer = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(10.0), Inches(user_y),
    Inches(3.0), Inches(database_y + layer_height - 0.3 - user_y)
)
api_layer.fill.solid()
api_layer.fill.fore_color.rgb = colors['layer_bg']
api_layer.line.color.rgb = colors['layer_border']
api_layer.line.width = Pt(1.5)

# API Layer Label
api_label = slide.shapes.add_textbox(Inches(10.2), Inches(user_y + 0.05), Inches(2), Inches(0.3))
api_label.text_frame.text = "5. API Layer"
api_label.text_frame.paragraphs[0].font.size = Pt(14)
api_label.text_frame.paragraphs[0].font.bold = True

# API Components with better spacing
api_components = [
    {"text": "Public API", "y": user_y + 0.45},
    {"text": "E-Saraban API", "y": user_y + 1.05},
    {"text": "Gazette Website API", "y": user_y + 1.65},
    {"text": "Standard APIs", "y": user_y + 2.25, "color": colors['backup']},
    {"text": "Security APIs", "y": user_y + 2.85, "color": RGBColor(211, 47, 47)},
    {"text": "Integration Testing", "y": user_y + 3.45, "color": colors['security']},
    {"text": "Penetration Testing", "y": user_y + 4.05, "color": colors['security']},
]

for comp in api_components:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.3), Inches(comp["y"]),
        Inches(2.6), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = comp.get("color", colors['api'])
    shape.text_frame.text = comp["text"]
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Save the presentation
prs.save('royal_gazette_architecture.pptx')
print("PowerPoint file 'royal_gazette_architecture.pptx' has been created successfully!")